import mysql.connector
from pptx import Presentation
from pptx.util import Pt
import os

# PowerPoint template path
template_path = "/home/fixroot/uploads/Impactslides.pptx"
output_dir = "/home/fixroot/output_presentations"
os.makedirs(output_dir, exist_ok=True)

# Database connection
conn = mysql.connector.connect(
    host="91.99.202.248",
    user="fixadmin",
    password="",
    database="FIXDB",
    port=3306
)
cursor = conn.cursor(dictionary=True)

# Fetch all rows
cursor.execute("""
SELECT
    Q1.GEMEENTE,
    HUISHOUDENS_FIXED,
    Q1.PERIODE,
    CONCAT(SUBSTR(Q1.PERIODE,6,2),' ',SUBSTR(Q1.PERIODE,1,4)) AS PERIODET,
    FORMAT(0.41025641*HUISHOUDENS_FIXED,0,'de_DE') AS CO2,
    FORMAT(20*HUISHOUDENS_FIXED,0,'de_DE') AS BOMEN,
    FORMAT(166*HUISHOUDENS_FIXED,0,'de_DE') AS GAS,
    FORMAT(407*HUISHOUDENS_FIXED,0,'de_DE') AS KWH,
    FORMAT(0.461538462*HUISHOUDENS_FIXED,0,'de_DE') AS TV,
    FORMAT(830*HUISHOUDENS_FIXED,0,'de_DE') AS douche,
    FORMAT(330*HUISHOUDENS_FIXED,0,'de_DE') AS SPAAR,
    QUOTE,
    concat(QUOTE_VAN,' uit ',Q1.GEMEENTE) as QUOTE_VAN
FROM FIXDB.HUISHOUDENS_GEMEENTE_PERIODE Q1
left join
FIXDB.QUOTES_VOOR_SLIDES Q2
on Q1.gemeente =Q2.gemeente
WHERE Q1.PERIODE='2025_Q2'
order by Q1.GEMEENTE
""")
rows = cursor.fetchall()

# Generate PowerPoint for each row
for row in rows:
    prs = Presentation(template_path)

    # Prepare replacement dictionary
    replacements = {
        "{{FIXBR}}": row["GEMEENTE"],
        "{{HH}}": row["HUISHOUDENS_FIXED"],
        "{{PERIODE}}": row["PERIODET"],
        "{{CO2}}": row["CO2"],
        "{{BOMEN}}": row["BOMEN"],
        "{{GAS}}": row["GAS"],
        "{{douche}}": row["douche"],
        "{{KWH}}": row["KWH"],
        "{{TV}}": row["TV"],
        "{{SPAAR}}": row["SPAAR"],
        "{{QUOTE}}": row["QUOTE"] if row.get("QUOTE") else "Niemand meer in de kou!"
    }

    # Replace text in slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Replace normal placeholders
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, replacement in replacements.items():
                            run.text = run.text.replace(placeholder, str(replacement))

                # Handle QUOTE_VAN separately with 10pt font
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and "{{QUOTE_VAN}}" in run.text:
                            replacement = row.get("QUOTE_VAN", "") or ""
                            run.text = run.text.replace("{{QUOTE_VAN}}", replacement)
                            run.font.size = Pt(10)          # 10pt lettergrootte

    # Save presentation
    output_path = os.path.join(output_dir, f"Impactslide_{row['GEMEENTE']}_{row['PERIODE']}.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")

# Cleanup
cursor.close()
conn.close()
