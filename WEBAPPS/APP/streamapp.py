import streamlit as st
import mysql.connector
from pptx import Presentation
from pptx.util import Pt
import io
from zipfile import ZipFile
import os
import subprocess
import tempfile

st.title("Generate Impact Slides ZIP")

# Function to build ZIP buffer (cached to avoid re-running on widget interaction)
@st.cache_data
def create_zip():
    # DB connection
    conn = mysql.connector.connect(
        host="91.99.202.248",
        user="fixadmin",
        password="",
        database="FIXDB",
        port=3306
    )
    cursor = conn.cursor(dictionary=True)
    cursor.execute("""
    SELECT
        Q1.GEMEENTE, HUISHOUDENS_FIXED, Q1.PERIODE,
        CONCAT(SUBSTR(Q1.PERIODE,6,2),' ',SUBSTR(Q1.PERIODE,1,4)) AS PERIODET,
        FORMAT(0.41025641*HUISHOUDENS_FIXED,0,'de_DE') AS CO2,
        FORMAT(20*HUISHOUDENS_FIXED,0,'de_DE') AS BOMEN,
        FORMAT(166*HUISHOUDENS_FIXED,0,'de_DE') AS GAS,
        FORMAT(407*HUISHOUDENS_FIXED,0,'de_DE') AS KWH,
        FORMAT(0.461538462*HUISHOUDENS_FIXED,0,'de_DE') AS TV,
        FORMAT(830*HUISHOUDENS_FIXED,0,'de_DE') AS douche,
        FORMAT(330*HUISHOUDENS_FIXED,0,'de_DE') AS SPAAR,
        QUOTE,
        concat(QUOTE_VAN,' uit ',Q1.GEMEENTE) AS QUOTE_VAN
    FROM FIXDB.HUISHOUDENS_GEMEENTE_PERIODE Q1
    LEFT JOIN FIXDB.QUOTES_VOOR_SLIDES Q2
      ON Q1.gemeente = Q2.gemeente
    WHERE Q1.PERIODE='2025_Q2'
    ORDER BY Q1.GEMEENTE
    """)
    rows = cursor.fetchall()
    cursor.close()
    conn.close()

    zip_buffer = io.BytesIO()
    with ZipFile(zip_buffer, "w") as zip_file:
        for row in rows:
            prs = Presentation("/home/fixroot/uploads/Impactslides.pptx")
            replacements = {
                "{{FIXBR}}": row["GEMEENTE"],
                "{{HH}}": row["HUISHOUDENS_FIXED"],
                "{{PERIODE}}": row["PERIODET"],
                "{{CO2}}": row["CO2"],
                "{{BOMEN}}": row["BOMEN"],
                "{{GAS}}": row["GAS"],
                "{{douche}}": row["douche"],
                "{{KWH}}": row["KWH"],
                "{{TV}}": row["TV"],
                "{{SPAAR}}": row["SPAAR"],
                "{{QUOTE}}": row["QUOTE"] or "Niemand meer in de kou!"
            }

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in paragraph.runs)
                            for ph, repl in replacements.items():
                                text = text.replace(ph, str(repl))
                            for i, run in enumerate(paragraph.runs):
                                run.text = text if i == 0 else ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "{{QUOTE_VAN}}" in run.text:
                                    rep = row.get("QUOTE_VAN", "") or ""
                                    run.text = run.text.replace("{{QUOTE_VAN}}", rep)
                                    run.font.size = Pt(10)

            with tempfile.TemporaryDirectory() as tmpdir:
                pptx_path = os.path.join(tmpdir, f"{row['GEMEENTE']}_{row['PERIODE']}.pptx")
                prs.save(pptx_path)
                with open(pptx_path, "rb") as f:
                    zip_file.writestr(os.path.basename(pptx_path), f.read())

                subprocess.run([
                    "libreoffice", "--headless",
                    "--convert-to", "png", "--outdir", tmpdir,
                    pptx_path
                ], check=True)

                for fn in sorted(os.listdir(tmpdir)):
                    if fn.lower().endswith(".png"):
                        with open(os.path.join(tmpdir, fn), "rb") as img:
                            zip_file.writestr(fn, img.read())

    zip_buffer.seek(0)
    return zip_buffer.getvalue()

# Button to generate ZIP
if st.button("Generate PPTX + PNG ZIP"):
    with st.spinner("Creating files..."):
        zip_data = create_zip()
    st.success("ZIP ready for download!")
    st.download_button(
        label="Download ZIP",
        data=zip_data,
        file_name="Impactslides.zip",
        mime="application/zip"
    )
